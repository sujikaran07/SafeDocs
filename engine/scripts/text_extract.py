from __future__ import annotations
from pathlib import Path

def extract_text(path: str) -> str:
    p = Path(path)
    suf = p.suffix.lower()

    try:
        if suf == ".pdf":
            from pdfminer.high_level import extract_text as pdf_extract
            return pdf_extract(str(p)) or ""
        elif suf in (".docx", ".pptx", ".xlsx"):
            # very light OOXML read (titles and text where easy)
            if suf == ".docx":
                from docx import Document
                doc = Document(str(p))
                return "\n".join([para.text for para in doc.paragraphs if para.text])
            elif suf == ".pptx":
                from pptx import Presentation
                pres = Presentation(str(p))
                chunks = []
                for slide in pres.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            chunks.append(shape.text)
                return "\n".join(chunks)
            elif suf == ".xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
                chunks = []
                for sh in wb.sheetnames:
                    ws = wb[sh]
                    for row in ws.iter_rows(values_only=True):
                        for cell in row:
                            if isinstance(cell, str) and cell.strip():
                                chunks.append(cell)
                return "\n".join(chunks)
        elif suf == ".rtf":
            # quick & dirty plain text strip
            import re
            txt = p.read_text(encoding="utf-8", errors="ignore")
            txt = re.sub(r"{\\[^}]+}", " ", txt)  # remove control groups
            txt = re.sub(r"\\[a-zA-Z]+-?\d*", " ", txt)
            return " ".join(txt.split())
    except Exception:
        return ""
    return ""
